"""
GitHub Presentation Generator
Creates a PPTX presentation with GitHub project information
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_github_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (GitHub-inspired)
    GITHUB_DARK = RGBColor(36, 41, 46)
    GITHUB_BLUE = RGBColor(0, 102, 204)
    ACCENT_GREEN = RGBColor(28, 211, 162)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GITHUB_DARK
    
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Library Management System"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "GitHub Project Overview & Developer Information"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Developer Information
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header
    header_shape = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Developer Information"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    info_items = [
        ("Name:", "Bhandari Prajwal"),
        ("User ID:", "25013467"),
        ("GitHub Username:", "prajwal3550"),
        ("Repository:", "library_management_system_py"),
        ("Profile URL:", "https://github.com/prajwal3550"),
    ]
    
    for label, value in info_items:
        p = text_frame.add_paragraph()
        p.text = f"{label} {value}"
        p.font.size = Pt(24)
        p.font.bold = True if label else False
        p.level = 0
        p.space_before = Pt(12)
    
    # Slide 3: Repository Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Repository Information"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    repo_info = [
        ("Repository Name:", "library_management_system_py"),
        ("Language:", "Python"),
        ("License:", "MIT License"),
        ("Repository URL:", "https://github.com/prajwal3550/library_management_system_py"),
        ("Repository Type:", "Public"),
        ("Default Branch:", "main"),
        ("Created:", "June 17, 2026"),
    ]
    
    for label, value in repo_info:
        p = text_frame.add_paragraph()
        p.text = f"{label}"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = value
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(8)
    
    # Slide 4: Project Features
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Project Features"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide4.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    features = [
        "Add new books to the system",
        "View all books in the library",
        "Search books by title",
        "Update book information",
        "Delete books from inventory",
        "Issue books to members",
        "Return books from members",
        "Persistent data storage using SQLite",
    ]
    
    for feature in features:
        p = text_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(20)
        p.level = 0
        p.space_before = Pt(10)
        p.font.color.rgb = GITHUB_BLUE
    
    # Slide 5: Technology Stack
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Technology Stack"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Two column layout
    left_box = slide5.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = "Backend"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GITHUB_DARK
    
    p = left_frame.add_paragraph()
    p.text = "Python 3"
    p.font.size = Pt(22)
    p.level = 1
    p.space_before = Pt(10)
    
    p = left_frame.add_paragraph()
    p.text = "SQLite3"
    p.font.size = Pt(22)
    p.level = 1
    
    p = left_frame.add_paragraph()
    p.text = "Command-line Interface"
    p.font.size = Pt(22)
    p.level = 1
    
    right_box = slide5.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = "Future Enhancements"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GITHUB_DARK
    
    enhancements = [
        "Graphical User Interface (GUI)",
        "Member Management System",
        "Due Date Tracking",
        "Fine Management",
        "Login System",
        "Export Reports",
    ]
    
    for enhancement in enhancements:
        p = right_frame.add_paragraph()
        p.text = enhancement
        p.font.size = Pt(16)
        p.level = 1
        p.space_before = Pt(6)
    
    # Slide 6: System Overview - Dashboard
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide6.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide6.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "System Overview - Dashboard"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "The Dashboard provides a quick overview of library statistics:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    dashboard_items = [
        ("Total Books:", "44 books in the collection"),
        ("Total Members:", "5 registered members"),
        ("Issued Books:", "0 books currently issued"),
        ("Available Books:", "44 books available for issuing"),
        ("Recent Transactions:", "Track of all book issues and returns"),
    ]
    
    for label, value in dashboard_items:
        p = text_frame.add_paragraph()
        p.text = f"{label} {value}"
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        
        if "Total" in label:
            p.font.color.rgb = GITHUB_BLUE
    
    # Slide 7: Book Management Module
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide7.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide7.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Book Management Module"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Book Details Available:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GITHUB_DARK
    
    book_details = [
        "Book ID - Unique identifier for each book",
        "Title - Name of the book",
        "Author - Name of the book author",
        "Category - Classification (Fiction, Self-Help, Technology, Sci-Fi, etc.)",
        "Quantity - Total copies available",
        "Available Count - Currently available books",
    ]
    
    for detail in book_details:
        p = text_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(18)
        p.space_before = Pt(8)
        p.font.color.rgb = GITHUB_BLUE
    
    # Slide 8: Member Management Module
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide8.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Member Management Module"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Member Information Tracked:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GITHUB_DARK
    
    member_info = [
        "Member ID - Unique identifier",
        "Full Name - Member's complete name",
        "Phone Number - Contact information",
        "Email Address - Email for communications",
    ]
    
    for info in member_info:
        p = text_frame.add_paragraph()
        p.text = info
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        p.font.color.rgb = GITHUB_BLUE
    
    p = text_frame.add_paragraph()
    p.text = "Current Members: 5"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(16)
    p.font.color.rgb = ACCENT_GREEN
    
    # Slide 9: Issue & Return Module
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide9.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide9.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Issue & Return Module"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide9.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Transaction Information:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GITHUB_DARK
    
    transactions = [
        "Transaction ID - Unique reference number",
        "Book Title - Name of issued/returned book",
        "Member Name - Who issued/returned the book",
        "Issue Date - Date of transaction",
        "Status - Current state (Issued/Returned)",
    ]
    
    for trans in transactions:
        p = text_frame.add_paragraph()
        p.text = trans
        p.font.size = Pt(18)
        p.space_before = Pt(10)
        p.font.color.rgb = GITHUB_BLUE
    
    p = text_frame.add_paragraph()
    p.text = "Current Active Issues: 0"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(16)
    p.font.color.rgb = RGBColor(220, 53, 69)
    
    # Slide 10: GitHub Collaboration
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide10.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide10.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "GitHub Repository Features"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide10.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    features_list = [
        "Public Repository - Open for collaboration",
        "MIT License - Free and open-source",
        "Issue Tracking - Report and manage bugs",
        "Pull Requests - Community contributions",
        "Wiki - Documentation support",
        "Project Board - Track development progress",
        "Downloads - Easy access to code",
    ]
    
    for feature in features_list:
        p = text_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(20)
        p.space_before = Pt(10)
        p.font.color.rgb = GITHUB_BLUE
    
    # Slide 11: Project Structure
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide11.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide11.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Project Structure"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide11.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "library-management-system/"
    p.font.size = Pt(18)
    p.font.color.rgb = GITHUB_DARK
    p.font.bold = True
    
    structure = [
        ("main.py", "Main application entry point"),
        ("library.db", "SQLite database file"),
        ("README.md", "Project documentation"),
        ("DOCUMENTATION.md", "Detailed documentation"),
        ("CONTRIBUTING.md", "Contribution guidelines"),
        ("LICENSE", "MIT License file"),
        (".gitignore", "Git ignore rules"),
    ]
    
    for file, desc in structure:
        p = text_frame.add_paragraph()
        p.text = f"├── {file}"
        p.font.size = Pt(14)
        p.font.color.rgb = GITHUB_BLUE
        p.level = 0
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.italic = True
        p.level = 1
    
    # Slide 12: Getting Started
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide12.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    header_shape = slide12.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GITHUB_DARK
    header_shape.line.color.rgb = GITHUB_DARK
    
    header_text = header_shape.text_frame
    header_text.text = "Getting Started"
    header_text.paragraphs[0].font.size = Pt(40)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = ACCENT_GREEN
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide12.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Installation Steps:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GITHUB_DARK
    
    steps = [
        "1. Clone the repository from GitHub",
        "2. Navigate to the project folder",
        "3. Ensure Python 3 is installed",
        "4. Run the application: python main.py",
        "5. Use the menu-driven interface to manage the library",
    ]
    
    for step in steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.space_before = Pt(10)
        p.font.color.rgb = GITHUB_BLUE
    
    # Slide 13: Contact & Resources
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide13.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GITHUB_DARK
    
    header_text_box = slide13.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.5))
    header_frame = header_text_box.text_frame
    p = header_frame.paragraphs[0]
    p.text = "Contact & Resources"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    info_box = slide13.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(4))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    contact_info = [
        ("Developer:", "Bhandari Prajwal (ID: 25013467)"),
        ("GitHub Profile:", "https://github.com/prajwal3550"),
        ("Repository:", "https://github.com/prajwal3550/library_management_system_py"),
        ("License:", "MIT License - Free to use and modify"),
    ]
    
    for label, value in contact_info:
        p = info_frame.add_paragraph()
        p.text = f"{label}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        p = info_frame.add_paragraph()
        p.text = value
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.level = 1
        p.space_after = Pt(12)
    
    # Save presentation
    output_path = "Library_Management_System_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    create_github_presentation()
